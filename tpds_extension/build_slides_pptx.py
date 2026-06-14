#!/usr/bin/env python3.11
"""Generate FitCachePP.pptx — a real PowerPoint deck (native charts + shapes).
Run: python3.11 tpds_extension/build_slides_pptx.py
Numbers reflect the 2026-06-14 TPDS revision (GNN regime-pinned, DRAM-tier proof, scale ladder)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

BLUE=RGBColor(0x25,0x63,0xEB); GREEN=RGBColor(0x16,0xA3,0x4A); AMBER=RGBColor(0xD9,0x77,0x06)
GRAY=RGBColor(0xB0,0xB6,0xBD); RED=RGBColor(0xDC,0x22,0x26); MUTED=RGBColor(0x5F,0x63,0x68)
FG=RGBColor(0x1B,0x1F,0x24); WHITE=RGBColor(0xFF,0xFF,0xFF); DGREEN=RGBColor(0x15,0x80,0x3D)
LBLUE=RGBColor(0xEE,0xF2,0xFF); LGREEN=RGBColor(0xDC,0xFC,0xE7); LAMBER=RGBColor(0xFE,0xF3,0xC7)
LRED=RGBColor(0xFE,0xE2,0xE2); CARD=RGBColor(0xF6,0xF8,0xFA); BORD=RGBColor(0xCB,0xD5,0xE1)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide():
    s=prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb=WHITE
    return s

def tbox(s,x,y,w,h,lines,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if isinstance(ln,tuple): ln=[ln]
        for (text,size,bold,color) in ln:
            r=p.add_run(); r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
            r.font.name="Calibri"
    return tb

def header(s,kicker,title):
    tbox(s,0.6,0.34,12.1,0.35,[[(kicker.upper(),12,True,BLUE)]])
    tbox(s,0.6,0.66,12.1,0.95,[[(title,27,True,FG)]])

def box(s,x,y,w,h,text,fill,fg=WHITE,size=12,bold=True,line=None):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.shadow.inherit=False
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    for m in ('margin_top','margin_bottom','margin_left','margin_right'): setattr(tf,m,Pt(3))
    for i,seg in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=seg; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=fg; r.font.name="Calibri"
    return sp

def rarrow(s,x,y,w,color=GRAY,h=0.24):
    sp=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def darrow(s,x,y,h,color=GRAY,w=0.24):
    sp=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False; return sp

def caption(s,x,y,w,h,text,size=11):
    tbox(s,x,y,w,h,[[(text,size,False,MUTED)]])

def chart(s,kind,x,y,w,h,cats,series,maxv=None,numfmt='0.0',ptcolors=None,sercolors=None,labels=True):
    cd=CategoryChartData(); cd.categories=cats
    for nm,vals in series: cd.add_series(nm,vals)
    gf=s.shapes.add_chart(kind,Inches(x),Inches(y),Inches(w),Inches(h),cd); ch=gf.chart
    ch.has_legend=len(series)>1
    if ch.has_legend:
        ch.legend.position=XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout=False; ch.legend.font.size=Pt(11)
    pl=ch.plots[0]; pl.gap_width=80
    if labels:
        pl.has_data_labels=True; dl=pl.data_labels; dl.number_format=numfmt; dl.number_format_is_linked=False
        dl.font.size=Pt(11); dl.font.bold=True
    if sercolors:
        for ser,clr in zip(ch.series,sercolors):
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb=clr
    if ptcolors:
        ser=ch.series[0]
        for i,clr in enumerate(ptcolors):
            ser.points[i].format.fill.solid(); ser.points[i].format.fill.fore_color.rgb=clr
    va=ch.value_axis
    if maxv is not None: va.maximum_scale=maxv; va.minimum_scale=0
    va.has_major_gridlines=True; va.tick_labels.font.size=Pt(10)
    ch.category_axis.tick_labels.font.size=Pt(11)
    for ax in (va,ch.category_axis): ax.format.line.color.rgb=BORD
    return ch

# ---------- 1 TITLE ----------
s=slide()
tbox(s,0.6,2.0,12.1,0.4,[[("TPDS EXTENSION · ADVISOR WALKTHROUGH",13,True,BLUE)]],PP_ALIGN.CENTER)
tbox(s,0.6,2.45,12.1,1.1,[[("FitCachePP",54,True,FG)]],PP_ALIGN.CENTER)
tbox(s,0.6,3.6,12.1,0.6,[[("Transparent local-storage caching for memory-mapped data loading on Frontier",20,False,MUTED)]],PP_ALIGN.CENTER)
b=box(s,2.4,4.5,8.5,1.3,"",CARD,line=BORD)
tf=b.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
for t,bd in [("Training jobs that read data with mmap (LLM shards, GNN features, vision tars) get data served ",False),
             ("transparently from a fast node-local cache",True),
             (" — with no application change. What the system does, then how much it helps (measured, checksum-verified).",False)]:
    r=p.add_run(); r.text=t; r.font.size=Pt(14); r.font.bold=bd; r.font.color.rgb=FG; r.font.name="Calibri"

# ---------- 2 PROBLEM ----------
s=slide(); header(s,"The problem","mmap-based loaders bypass the cache")
tbox(s,0.6,1.55,12.1,0.7,[[("Modern loaders read by memory-mapping a file and page-faulting — they never call read(), so they slip past a read()/pread() cache straight to Lustre.",14,False,FG)]])
box(s,0.7,2.7,2.6,0.8,"app: read() / pread()",LBLUE,FG,12,False,BORD)
rarrow(s,3.45,2.95,0.7,GRAY)
box(s,4.3,2.7,2.4,0.8,"FitCachePP hook",LBLUE,FG,12,False,BORD)
rarrow(s,6.85,2.95,0.7,GREEN)
box(s,7.7,2.7,3.4,0.8,"served from local cache",LGREEN,DGREEN,13,True,GREEN)
box(s,0.7,4.0,2.6,0.8,"app: mmap() page-fault",LBLUE,FG,12,False,BORD)
rarrow(s,3.45,4.25,3.4,RED)
tbox(s,3.6,3.95,3.2,0.3,[[("bypasses the cache",10,True,RED)]],PP_ALIGN.CENTER)
box(s,7.7,4.0,3.4,0.8,"straight to Lustre PFS",LRED,RED,13,True,RED)
caption(s,0.7,5.2,11.6,1.0,"On the workloads we target the cache measured 0 cache hits — every byte came from Lustre. Lustre streams ~2 GB/s sequentially, but sparse-random 4 KB page-faults crawl at ~3–4 MB/s — which is why the access pattern decides the win.")

# ---------- 3 OVERVIEW ----------
s=slide(); header(s,"What we do","Two contributions + one supporting mechanism")
cards=[("1","Transparent mmap cache path","Intercept mmap; on a warm hit hand back a direct map of the cached copy on local storage — deterministic, no network round-trip.",BLUE),
       ("2","Access-aware placement","Route each file's cached copy to the node that reads it. Removes cross-node misses; partitioned access stores one copy, not N.",BLUE),
       ("Supporting","Sidecar recovery","A tiny .meta per file lets a restarted cache rebuild its state in seconds — the cache survives a restart.",MUTED)]
for i,(tag,ttl,body,tagc) in enumerate(cards):
    x=0.7+i*4.05
    c=box(s,x,2.2,3.7,3.4,"",CARD,line=BORD)
    tg=box(s,x+0.25,2.45,1.7 if tag!='Supporting' else 1.9,0.45,tag,tagc,WHITE,11,True)
    tbox(s,x+0.25,3.1,3.2,0.7,[[(ttl,16,True,FG)]])
    tbox(s,x+0.25,3.9,3.25,1.6,[[(body,12.5,False,FG)]])

# ---------- 4 MECH 1 FLOW ----------
s=slide(); header(s,"What we do — 1 / 3","Transparent mmap cache path")
tbox(s,0.6,1.5,12.1,0.7,[[("On a warm hit the client recomputes the cached path (same hash as the server), checks it's complete, and maps the local file — kernel faults from local storage (DRAM/NVMe), no network round-trip.",13.5,False,FG)]])
box(s,5.2,2.45,2.9,0.7,"app calls mmap(file)",LBLUE,FG,13,True,BORD)
darrow(s,6.55,3.2,0.4,GRAY)
box(s,3.6,3.7,6.1,0.95,"client resolver  (no network)\nrecompute cached path · same hash as server · check size complete",CARD,FG,12,True,BORD)
darrow(s,5.0,4.75,0.35,GREEN); darrow(s,8.3,4.75,0.35,AMBER)
box(s,2.2,5.2,4.0,1.15,"WARM HIT\nmmap the local cached file → kernel faults from the local tier (DRAM/NVMe)",LGREEN,DGREEN,12,True,GREEN)
box(s,7.1,5.2,4.0,1.15,"COLD MISS\nnative PFS mmap (never slower than baseline) → server promotes in background → next mmap = warm",LAMBER,RGBColor(0xB4,0x53,0x09),12,True,RGBColor(0xFC,0xD3,0x4D))
caption(s,0.6,6.55,12.1,0.5,"The warm-hit path has no RPC: the client finds the cached file by computing its path, then maps it directly.")

# ---------- 5 MECH 1 WHY NO RPC ----------
s=slide(); header(s,"What we do — 1 / 3","Why map a local file instead of fetching over the network?")
tbox(s,0.6,1.5,12.1,0.5,[[("Our first design fetched data over Mercury RPC at mmap time — the slowest option. So the read path now touches no network at all.",13.5,False,FG)]])
chart(s,XL_CHART_TYPE.BAR_CLUSTERED,1.2,2.2,10.8,3.4,
      ["old: fetch over RPC","Lustre PFS (baseline)","direct local mmap"],
      [("GB/s",(0.9,2.0,2.3))],maxv=2.6,numfmt='0.0',ptcolors=[GRAY,GRAY,BLUE])
caption(s,0.6,5.9,12.1,0.8,"Single-stream bandwidth (GB/s). RPC ≈0.9 and Lustre ≈2.0 are from design notes; the local-tier bar is an estimate showing it is faster than Lustre. The point is the ordering — the old RPC fetch was slower than Lustre itself, so we dropped it.")

# ---------- 6 MECH 2 SCATTER ----------
s=slide(); header(s,"What we do — 2 / 3","Access-aware placement")
tbox(s,0.6,1.5,12.1,0.6,[[("A warm hit needs the file on the same node as the rank reading it. Default hashing scatters a node's files across the cluster, so it misses ~half its own data. We route each cached copy to its own node.",13.5,False,FG)]])
box(s,0.9,2.5,5.5,0.6,"Default hash — node 0's 40 files",WHITE,FG,13,True,BORD)
box(s,1.2,3.25,2.4,0.7,"18 on this node",GREEN,WHITE,13,True)
box(s,3.7,3.25,2.5,0.7,"22 scattered away",RED,WHITE,13,True)
tbox(s,0.9,4.1,5.5,0.4,[[("→ 22 fall back to Lustre",13,True,RED)]],PP_ALIGN.CENTER)
box(s,6.9,2.5,5.5,0.6,"Node-local — node 0's 40 files",WHITE,FG,13,True,BORD)
box(s,7.2,3.25,4.9,0.7,"all 40 on this node · 0 remote",GREEN,WHITE,13,True)
tbox(s,6.9,4.1,5.5,0.4,[[("→ every read is a warm hit",13,True,DGREEN)]],PP_ALIGN.CENTER)
caption(s,0.6,4.9,12.1,0.5,"Same data, two placements. Node-local routing turns 22 cross-node misses into 0 (measured: Megatron, N=2).")

# ---------- 7 MECH 2 PARTITIONED ----------
s=slide(); header(s,"What we do — 2 / 3","Partitioned access stores the corpus once, not N times")
chart(s,XL_CHART_TYPE.COLUMN_CLUSTERED,0.9,2.0,5.4,3.6,["hash","node-local"],
      [("cross-node misses/node",(22,0))],maxv=24,numfmt='0',ptcolors=[RED,GREEN])
tbox(s,0.9,1.65,5.4,0.3,[[("Cross-node misses per node",13,True,FG)]])
chart(s,XL_CHART_TYPE.COLUMN_CLUSTERED,7.0,2.0,5.4,3.6,["shared","partitioned"],
      [("copies stored",(2.0,1.0))],maxv=2.4,numfmt='0.0"×"',ptcolors=[AMBER,GREEN])
tbox(s,7.0,1.65,5.4,0.3,[[("Copies of the corpus stored",13,True,FG)]])
b=box(s,0.9,5.85,11.5,0.9,"",LGREEN,line=GREEN); tf=b.text_frame; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
for t,bd in [("Both get full local hits (0 misses). ",False),("Partitioned placement gives the same speed at half the disk",True),(" — at 16 nodes, the difference between 1 copy of the dataset and 16.",False)]:
    r=p.add_run(); r.text=t; r.font.size=Pt(13); r.font.bold=bd; r.font.color.rgb=FG; r.font.name="Calibri"

# ---------- 8 MECH 3 SIDECAR ----------
s=slide(); header(s,"What we do — 3 / 3","Sidecar recovery")
tbox(s,0.6,1.5,12.1,0.5,[[("Each cached file gets a tiny .meta sidecar. On restart the cache rebuilds its state from those sidecars in seconds — it survives the restart instead of re-fetching from the filesystem.",13.5,False,FG)]])
box(s,0.8,2.5,2.9,0.8,"promote file →\nwrite .meta sidecar",LBLUE,FG,12,False,BORD)
rarrow(s,3.85,2.75,0.6,GRAY)
box(s,4.6,2.5,2.3,0.8,"server restart",LAMBER,RGBColor(0xB4,0x53,0x09),13,True,RGBColor(0xFC,0xD3,0x4D))
rarrow(s,7.05,2.75,0.6,GRAY)
box(s,7.8,2.5,3.0,0.8,"scan tiers, rebuild\ncache map (~5 s)",LBLUE,FG,12,False,BORD)
rarrow(s,10.95,2.75,0.55,GREEN)
box(s,11.65,2.5,1.1,0.8,"warm",LGREEN,DGREEN,12,True,GREEN)
chart(s,XL_CHART_TYPE.LINE_MARKERS,2.0,3.7,9.0,2.7,["32 objects","256 objects","996 objects"],
      [("restore (s)",(6,5,5))],maxv=8,numfmt='0" s"',sercolors=[BLUE])
caption(s,0.6,6.5,12.1,0.6,"Metadata-restore-only sweep: restore time is independent of object count (~1 KB each). Flat → scales to 1000s of objects. (A full DINOv2 recovery, which re-checks data files, takes ~8 s — a heavier path.)")

# ---------- 9 RESULTS dense ----------
s=slide(); header(s,"How much faster — 1 / 3","Dense-sequential: a modest, stable win")
tbox(s,0.6,1.5,12.1,0.5,[[("Warm-hit speedup of local-cache mmap vs native PFS mmap. Magnitude is reported per regime; every number passed a bit-for-bit checksum-equality check.",13.5,False,FG)]])
chart(s,XL_CHART_TYPE.BAR_CLUSTERED,1.2,2.2,10.8,3.3,["Megatron (LLM)","DINOv2 (vision)"],
      [("warm-hit speedup",(1.6,1.95))],maxv=2.5,numfmt='0.00"×"',ptcolors=[BLUE,GREEN])
caption(s,0.6,5.7,12.1,0.9,"Bars show the high end of each range (Megatron 1.2–1.6×, DINOv2 1.6–1.95×). A modest but stable warm-hit win, flat across N=1–16, because Frontier's Lustre is already fast at streaming. The sparse-random GNN workload is a different story — next.")

# ---------- 10 RESULTS GNN regime ----------
s=slide(); header(s,"How much faster — 2 / 3","GNN sparse-random: the win depends on the baseline regime")
chart(s,XL_CHART_TYPE.BAR_CLUSTERED,1.2,1.8,10.8,3.4,
      ["page-cache (best)","NVMe-bound warm","cold Lustre native"],
      [("epoch (s)",(5,82.7,2195))],maxv=2600,numfmt='0" s"',ptcolors=[GREEN,BLUE,GRAY])
b=box(s,0.9,5.5,11.5,1.25,"",CARD,line=BORD); tf=b.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
for t,bd,cl in [("2026-06-14 regime-pinned restatement (replaces the earlier 18–33×, which was regime-contaminated). ",True,FG),
                ("The claim is the NVMe-bound 82.7 s: ~6× faster than the most favorable Lustre native (510 s), and 22–30× faster than cold Lustre (1870–2457 s). ",False,FG),
                ("Source gnn_regime_pinned/, checksum-gated.",False,MUTED)]:
    r=p.add_run(); r.text=t; r.font.size=Pt(12.5); r.font.bold=bd; r.font.color.rgb=cl; r.font.name="Calibri"

# ---------- 11 RESULTS scale ladder ----------
s=slide(); header(s,"How much faster — 3 / 3","Warm stays flat as the job scales")
chart(s,XL_CHART_TYPE.COLUMN_CLUSTERED,0.9,1.8,11.5,3.5,["32 nodes","64 nodes","128 nodes"],
      [("FitCachePP warm",(176,176,176)),("Native cold (Lustre)",(2568,944,1066))],
      maxv=2700,numfmt='0',sercolors=[BLUE,GRAY])
tbox(s,0.9,5.35,11.5,0.5,[[("Warm flat 175–178 s across 32/64/128 nodes; speedup 14.7× / 5.3× / 6.0× varies only with the native baseline's Lustre contention.",12.5,False,FG)]])
b=box(s,0.9,5.95,11.5,0.85,"",LGREEN,line=GREEN); tf=b.text_frame; p=tf.paragraphs[0]
for t,bd in [("The warm path is multi-tier, not NVMe-only: ",True),("a controlled run served 48 warm hits from the DRAM (/tmp tmpfs) tier, bit-identical and at parity with NVMe.",False)]:
    r=p.add_run(); r.text=t; r.font.size=Pt(12.5); r.font.bold=bd; r.font.color.rgb=FG; r.font.name="Calibri"

# ---------- 12 HONEST FRAMING ----------
s=slide(); header(s,"Honest framing","What the numbers do and don't claim")
items=[("This accelerates data loading, not GPU compute","— we speed up the I/O the loader does."),
       ("Warm-hit runs use a pre-staged cache","— we also report the amortized view including that staging cost (break-even 1–4 epochs)."),
       ("Magnitude is reported per regime","— GNN ~6× vs favorable Lustre / 22–30× vs cold; dense-sequential ~1.2–1.95×."),
       ("Tier-general, served from NVMe here","— also warm-hits DRAM/tmpfs; working sets (640 GiB–902 GB/node) exceed node RAM, so they land on NVMe, keeping the baseline from cheating via the page cache.")]
y=2.0
for bold,rest in items:
    bx=box(s,0.8,y,11.7,1.0,"",RGBColor(0xFF,0xFB,0xEB),line=RGBColor(0xFD,0xE6,0x8A))
    tf=bx.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    r=p.add_run(); r.text=bold; r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=FG; r.font.name="Calibri"
    r=p.add_run(); r.text=" "+rest; r.font.size=Pt(14); r.font.bold=False; r.font.color.rgb=FG; r.font.name="Calibri"
    y+=1.15

# ---------- 13 STATUS ----------
s=slide(); header(s,"Status & next steps","Where it stands")
rows=[("Component","State","Evidence"),
      ("Transparent mmap cache path","Done","3 workload classes, N=1–16; GNN regime-pinned (~6× / 22–30×)"),
      ("Multi-tier warm hit (DRAM+NVMe)","Done","48 tmpfs hits, bit-identical, parity with NVMe"),
      ("Warm-flat at scale (Megatron)","32–128N (256N pending)","warm flat 175–178 s; native swings with contention"),
      ("Access-aware placement","Done","0 cross-node misses; partitioned = 1× storage"),
      ("Sidecar recovery","Done","flat ~5 s restore to ~1000 objects"),
      ("Cross-job sharing","Future work","preliminary; known correctness defect; not a claim")]
tb=s.shapes.add_table(len(rows),3,Inches(0.7),Inches(1.7),Inches(11.9),Inches(3.0)).table
tb.columns[0].width=Inches(4.3); tb.columns[1].width=Inches(2.9); tb.columns[2].width=Inches(4.7)
for ci,_ in enumerate(rows[0]):
    for ri,row in enumerate(rows):
        cell=tb.cell(ri,ci); cell.margin_left=Pt(6); cell.margin_top=Pt(3); cell.margin_bottom=Pt(3)
        para=cell.text_frame.paragraphs[0]; r=para.add_run(); r.text=row[ci]
        r.font.size=Pt(12); r.font.name="Calibri"; r.font.bold=(ri==0); r.font.color.rgb=(WHITE if ri==0 else FG)
        cell.fill.solid(); cell.fill.fore_color.rgb=(BLUE if ri==0 else (WHITE if ri%2 else CARD))
tbox(s,0.7,5.0,11.9,1.2,[[("Next: ",13,True,FG),("an organic-promotion run (no manual pre-stage) + a manual-staging baseline; a contended-filesystem scenario where the win grows; a hash-vs-node-local timing number; and one end-to-end training number tying loader speedup to job wall-clock.",13,False,FG)]])

import os
out=os.path.join(os.path.dirname(__file__),"FitCachePP.pptx")
prs.save(out)
print("WROTE",out,"slides:",len(prs.slides._sldIdLst))
